from pptx import Presentation

source1 = "../../Source/Baseball Rules!.pptx"
source2 = "../../Source/HiTech GmbH Vision.pptx"
source3 = "../../Source/Präsentation_LF02-v2.pptx"

def analyze_single_slide(filepath, layout_index):
    # Load your presentation template
    prs = Presentation(filepath)
    # Choose the slide layout index you're interested in
    # Slide layouts are indexed starting from 0

    # Get the layout
    slide_layout = prs.slide_layouts[layout_index]

    # Print placeholder details for the chosen layout
    # for shape in slide_layout.shapes:
    #     print(f"Shape Type: {shape.shape_type}, Shape ID: {shape.shape_id}")
    print(f"Details for layout {layout_index}: {slide_layout.name}")

    for placeholder in slide_layout.placeholders:
        if hasattr(placeholder, "text"):
            print(placeholder.text)
        print(f"Placeholder index: {placeholder.placeholder_format.idx}, Type: {placeholder.placeholder_format.type}, "
              f"Name: '{placeholder.name}' dimensions: left: {placeholder.left}, top: {placeholder.top},"
              f"width: {placeholder.width}, height: {placeholder.height} ")

    slides_text = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text.append("\n".join(slide_text))

    for i, text in enumerate(slides_text):
        print(f"Slide {i + 1}:\n{text}\n")



def analyze_presentation(filepath):
    print("\n"f"Presentation: {filepath}")
    presentation = Presentation(filepath)
    for slide_index, slide in enumerate(presentation.slides):
        # Find the layout index by matching the layout object
        layout_index = next((i for i, layout in enumerate(presentation.slide_layouts) if layout == slide.slide_layout), None)
        if layout_index is None:
            print(f"Layout for slide {slide_index + 1} not found.")
            continue

        print(f"\nAnalyzing slide {slide_index + 1} :")
        analyze_single_slide(filepath, layout_index)

        # Check for placeholders on the slide
        for placeholder in slide.placeholders:
            print(f"Slide {slide_index + 1} has placeholder with index {placeholder.placeholder_format.idx}")



# Beispielaufruf
analyze_presentation(source3)


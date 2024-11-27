from pptx.util import Inches


def create_multi_slide_table(presentation, data, max_rows_per_slide=4):
    start_row = 0
    while start_row < len(data["table_content"]):
        start_row = create_table_slide(presentation, data, start_row, max_rows_per_slide)


def create_table_slide(presentation, data, start_row, max_rows_per_slide):
    table_slide_layout = presentation.slide_layouts[16]
    table_slide = presentation.slides.add_slide(table_slide_layout)
    title = table_slide.placeholders[0]
    title.text = data["headline"]
    subtitle = table_slide.placeholders[36]
    subtitle.text = data["sub_header"]
    # number of rows
    num_rows = min(len(data["table_content"]) - start_row, max_rows_per_slide) + 1  # +1 for table head
    cols = len(data["table_content"][0])
    # dimensions of table
    left = Inches(0.37)
    top = Inches(1.82)
    width = Inches(12.6)
    height = Inches(4.5)
    graphic_frame = table_slide.shapes.add_table(num_rows, cols, left, top, width, height)
    table = graphic_frame.table
    # table header height
    table.rows[0].height = Inches(0.5)
    # table header
    key_names = list(data["table_content"][0].keys())
    for i, key in enumerate(key_names):
        table.cell(0, i).text = key
    # table body
    for i in range(1, num_rows):
        entry = data["table_content"][start_row + i - 1]
        for j, key in enumerate(key_names):
            table.cell(i, j).text = entry[key]
    # returns rest rows if number of rows higher than maximal rows
    return start_row + num_rows - 1  # -1 for table head

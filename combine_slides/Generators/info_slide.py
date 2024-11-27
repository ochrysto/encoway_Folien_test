import tempfile
import cv2
from pptx.util import Inches
from combine_slides.Util import edit_image


def create_slide_9(presentation, data):
    slide_layout_9 = presentation.slide_layouts[9]
    slide = presentation.slides.add_slide(slide_layout_9)
    title = slide.shapes.placeholders[0]
    title.text = data["headline"]
    subtitle = slide.shapes.placeholders[36]
    subtitle.text = data["sub_header"]

    body_left = slide.shapes.placeholders[11]
    text_frame = body_left.text_frame

    for point in data["bullet_points"]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = point

    left = Inches(6.85)
    top = Inches(1.75)
    width = Inches(5)

    image = data["img_url"]
    img_format = edit_image.check_orientation(image)

    if img_format == "Portrait":
        cropped_image = edit_image.crop_image(image)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as temp_file:
            temp_file_path = temp_file.name
            cv2.imwrite(temp_file_path, cropped_image)
        slide.shapes.add_picture(temp_file_path, left, top, width)
    else:
        slide.shapes.add_picture(data["img_url"], left, top, width)


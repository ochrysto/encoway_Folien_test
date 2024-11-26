from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches


def create_chart_slide(presentation, data):
    chart_slide_layout = presentation.slide_layouts[17]
    chart_slide = presentation.slides.add_slide(chart_slide_layout)
    title = chart_slide.placeholders[0]
    title.text = data["headline"]
    subtitle = chart_slide.placeholders[36]
    subtitle.text = data["sub_header"]

    chart_data = ChartData()
    chart_data.categories = data["categories"]
    chart_title = data["series"][0]
    chart_values = data["series"][1:]
    chart_data.add_series(chart_title, chart_values)
    left = Inches(0.37)
    top = Inches(1.82)
    width = Inches(12.6)
    height = Inches(4.8)
    chart = chart_slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END



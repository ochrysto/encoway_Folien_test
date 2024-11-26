from pptx.util import Inches


def create_table_slide(presentation, data):
    table_slide_layout = presentation.slide_layouts[16]
    table_slide = presentation.slides.add_slide(table_slide_layout)
    title = table_slide.placeholders[0]
    title.text = data["headline"]
    subtitle = table_slide.placeholders[36]
    subtitle.text = data["sub_header"]

    rows = len(data["table_content"]) + 1
    cols = 4
    left = Inches(0.37)
    top = Inches(1.82)
    width = Inches(12.6)
    height = Inches(4.5)
    graphic_frame = table_slide.shapes.add_table(rows, cols, left, top, width, height)
    table = graphic_frame.table
    # table head
    key_names = list(data["table_content"][0].keys())

    table.cell(0, 0).text = key_names[0]
    table.cell(0, 1).text = key_names[1]
    table.cell(0, 2).text = key_names[2]
    table.cell(0, 3).text = key_names[3]
    #table body
    for i, entry in enumerate(data["table_content"], start=1):
        table.cell(i, 0).text = entry['name']
        table.cell(i, 1).text = entry['phone']
        table.cell(i, 2).text = entry['email']
        table.cell(i, 3).text = entry['date']


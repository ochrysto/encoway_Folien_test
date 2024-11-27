import unittest
from pptx import Presentation


def create_title_slide(presentation, data):
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = data["headline"]
    subtitle = title_slide.placeholders[10]
    subtitle.text = f"{data['name']} | {data['date']}"


class TestCreateTitleSlide(unittest.TestCase):
    def test_create_title_slide(self):
        # Arrange
        presentation = Presentation()
        data = {
            "headline": "Test Headline",
            "name": "Test Name",
            "date": "2024-11-27"
        }

        # Act
        create_title_slide(presentation, data)

        # Assert
        slide = presentation.slides[0]
        title = slide.shapes.title
        subtitle = slide.placeholders[10]

        self.assertEqual(title.text, "Test Headline")
        self.assertEqual(subtitle.text, "Test Name | 2024-11-27")


if __name__ == '__main__':
    unittest.main()


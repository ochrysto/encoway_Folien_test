import os
import json

import pptx
from pptx import Presentation

# Set the input folder path
input_folder = '../../Source'

# Create an empty list to store slide data
slides = []

# Loop through all PowerPoint files in the input folder
for filename in os.listdir(input_folder):
    if filename.endswith('.pptx'):
        # Load the PowerPoint file using pptx library
        prs = Presentation(os.path.join(input_folder, filename))

        # Iterate through each slide in the presentation
        for slide in prs.slides:
            # Check the slide type (title, info, chart, or table)
            if slide.slide_layout == 1:  # Title slide
                slide_data = {
                    "slide_format": "title",
                    "headline": slide.placeholders[0].text,
                    "sub_header": "",
                    "name": "",
                    "date": ""
                }
            elif slide.slide_layout in [2, 3]:  # Info slide with bulletpoints
                bullet_points = []
                for shape in slide.shapes:
                    if isinstance(shape, pptx.enum.shape.Shape):
                        text = ''
                        for paragraph in shape.text:
                            text += paragraph.text.strip()
                        bullet_points.append(text)
                slide_data = {
                    "slide_format": "info",
                    "headline": slide.placeholders[0].text,
                    "sub_header": "",
                    "img_url": "",  # Assume no image
                    "bullet_points": bullet_points
                }
            elif slide.slide_layout == 4:  # Table slide
                table_content = []
                for row in slide.tables[0].rows:
                    table_row = []
                    for cell in row.cells:
                        table_row.append(cell.text.strip())
                    table_content.append(table_row)
                slide_data = {
                    "slide_format": "table",
                    "headline": slide.placeholders[0].text,
                    "sub_header": "",
                    "table_content": table_content
                }
            elif slide.slide_layout == 5:  # Chart slide (assuming a simple bar chart for now)
                categories = []
                series = []
                for shape in slide.shapes:
                    if isinstance(shape, pptx.enum.shape.Shape):
                        text = ''
                        for paragraph in shape.text:
                            text += paragraph.text.strip()
                        if not categories:
                            categories.append(text)
                        else:
                            series.append(text)
                slide_data = {
                    "slide_format": "chart",
                    "chart_type": "simple_bar",
                    "headline": slide.placeholders[0].text,
                    "sub_header": "",
                    "categories": categories,
                    "series": series
                }
            else:  # Unknown slide type (e.g., end slide)
                slide_data = {
                    "slide_format": "end",
                    "headline": slide.placeholders[0].text
                }

            # Add the slide data to the overall slides list
            slides.append(slide_data)

with open('output.json', 'w') as f:
    json.dump({
        "template": "../../Material/encoway PowerPoint-12-2023-Master.pptx",
        "file_name": "Awesome-Beginning",
        "title": "Awesome Beginning",
        "slides": slides
    }, f, indent=4)
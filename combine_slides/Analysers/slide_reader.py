import json
import os
from typing import Any

from pptx import Presentation

input_folder = "../../Source"
output_json = "presentations.json"


def pptx_to_json(input_folder, output_json):
    presentations = []

    for filename in os.listdir(input_folder):
        if filename.endswith(".pptx"):
            filepath = os.path.join(input_folder, filename)
            presentation = Presentation(filepath)
            slides = []

            for slide in presentation.slides:
                slide_data: dict[str | Any, str | list[Any] | list[dict[Any, Any]] | Any] = {
                    "slide_format": "unknown",
                    "headline": "",
                    "sub_header": "",
                    "name": "",
                    "date": "",
                    "img_url": "",
                    "bullet_points": [],
                    "table_content": [],
                    "chart_type": "",
                    "categories": [],
                    "series": []
                }

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        if "title" in shape.name.lower():
                            slide_data["headline"] = text
                        elif "subtitle" in shape.name.lower():
                            slide_data["sub_header"] = text
                        else:
                            slide_data["bullet_points"].append(text)
                    elif shape.has_table:
                        table_content = []
                        for row in shape.table.rows:
                            row_data = {}
                            for cell in row.cells:
                                row_data[cell.text] = cell.text
                            table_content.append(row_data)
                        slide_data["table_content"] = table_content
                    # TODO weitere Bedingungen hinzufügen

                slides.append(slide_data)

            presentations.append({
                "file_name": filename,
                "slides": slides
            })

    with open(output_json, 'w') as f:
        json.dump(presentations, f, indent=4)

pptx_to_json(input_folder, output_json)
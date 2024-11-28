from tempfile import NamedTemporaryFile
import os
from pptx import Presentation

def extract_slide(presentation_file, slide_number, output_file):
    """
    Extrahiert eine bestimmte Folie aus einer PowerPoint-Präsentation und speichert sie in einer neuen Präsentation.

    :param presentation_file: Pfad zur Quellpräsentation.
    :param slide_number: Nummer der Folie, die extrahiert werden soll (1-basiert).
    :param output_file: Pfad zur neuen Präsentation, die die extrahierte Folie enthält.
    """
    # Öffne die bestehende Präsentation
    prs = Presentation(presentation_file)

    # PPTX-Folien sind nullbasiert, daher -1
    slide_index = slide_number - 1
    if slide_index < 0 or slide_index >= len(prs.slides):
        print(f"Warnung: Folie {slide_number} existiert nicht in {presentation_file}.")
        return

    # Erstelle eine neue leere Präsentation
    new_presentation = Presentation()

    # Hole die Folie und füge sie zur neuen Präsentation hinzu
    slide = prs.slides[slide_index]
    copy_slide(new_presentation, slide)

    # Speichere die neue Präsentation
    new_presentation.save(output_file)
    print(f"Neue Präsentation gespeichert unter: {output_file}")

def copy_slide(presentation, slide):
    """
    Kopiert eine Folie aus einer Präsentation in eine andere.

    :param presentation: Ziel-Präsentation.
    :param slide: Die Folie, die kopiert werden soll.
    """
    # Füge eine leere Folie mit einem Standardlayout hinzu
    layout = presentation.slide_layouts[0]  # Verwende ein Standardlayout
    new_slide = presentation.slides.add_slide(layout)

    # Kopiere Inhalte der Folie
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        if shape.shape_type == 1:  # Textbox
            new_shape = new_slide.shapes.add_shape(
                shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height
            )
            if shape.text:
                new_shape.text = shape.text
        elif shape.shape_type == 13:  # Bild
            with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile:
                tmpfile.write(shape.image.blob)
                tmpfile.flush()
                new_slide.shapes.add_picture(
                    tmpfile.name, shape.left, shape.top, shape.width, shape.height
                )
            os.remove(tmpfile.name)
        # Weitere Formtypen können bei Bedarf hinzugefügt werden

if __name__ == "__main__":
    # Beispiel: Extrahieren von Folie 3 aus einer Präsentation
    presentation_file = "../../Source/Präsentation_LF02-v2.pptx"
    slide_number = 2
    output_file = "../../Outcome/extracted_slide.pptx"
    extract_slide(presentation_file, slide_number, output_file)

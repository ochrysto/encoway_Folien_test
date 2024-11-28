import os
from tempfile import NamedTemporaryFile
from pptx import Presentation

def list_files_of_type(directory, file_extension):
    files = []
    for filename in os.listdir(directory):
        if filename.endswith(file_extension):
            files.append(filename)
    return files


def extract_and_combine_slides(presentation_files, slide_indices, source_path):
    # Create a new presentation for the combined slides
    combined_presentation = Presentation()

    for file, indices in zip(presentation_files, slide_indices):
        filepath = os.path.join(source_path, file)
        #filepath = f"{source_path}/{file}"
        if not os.path.exists(filepath):
            print(f"File not found: {filepath}")
            continue
        presentation = Presentation(filepath)

        for index in indices:
            if index < len(presentation.slides):
                slide = presentation.slides[index]
                slide_layout = combined_presentation.slide_layouts[5]
                new_slide = combined_presentation.slides.add_slide(slide_layout)

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        new_shape = new_slide.shapes.add_shape(
                            shape.shape_type, shape.left, shape.top, shape.width, shape.height
                        )
                        new_shape.text = shape.text_frame.text
                    elif shape.shape_type == 19:  # Placeholder
                        new_placeholder = new_slide.placeholders[shape.placeholder_format.idx]
                        new_placeholder.text = shape.text_frame.text
                    elif shape.shape_type == 13:  # Picture
                        with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile:
                            tmpfile.write(shape.image.blob)
                            tmpfile.flush()
                            new_slide.shapes.add_picture(
                                tmpfile.name, shape.left, shape.top, shape.width, shape.height
                            )
                        os.remove(tmpfile.name)
                    elif shape.has_table:
                        new_table = new_slide.shapes.add_table(
                            shape.table.rows, shape.table.columns, shape.left, shape.top, shape.width, shape.height
                        )
                        for r in range(shape.table.rows):
                            for c in range(shape.table.columns):
                                new_table.table.cell(r, c).text = shape.table.cell(r, c).text
                    # Add other shape types as needed

    combined_presentation.save(output_file)

source = "../../Source"
# Example usage
presentation_files = list_files_of_type(source , ".pptx")
print(presentation_files)
# Indices of slides to extract from each presentation
slide_indices = [[0, 1, 2], [1, 2], []]
output_file = "../../Outcome/combined_presentation.pptx"

extract_and_combine_slides(presentation_files, slide_indices, source)

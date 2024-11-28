import os
from tempfile import NamedTemporaryFile

from pptx import Presentation


def extract_and_merge_slides(slide_ranges, output_file):
    """
    Kombiniert Folien aus mehreren PowerPoint-Präsentationen basierend auf angegebenen Bereichen.

    :param slide_ranges: Liste von Tupeln (dateiname, [foliennummern]), z. B. [("alpha.pptx", [1, 2, 3]), ("beta.pptx", [2, 4])]
    :param output_file: Dateiname für die neue Präsentation.
    """
    # Erstelle eine neue leere Präsentation
    new_presentation = Presentation()

    for file, slide_numbers in slide_ranges:
        # Öffne die bestehende Präsentation
        prs = Presentation(file)

        for slide_number in slide_numbers:
            # PPTX-Folien sind nullbasiert, daher -1
            slide_index = slide_number - 1
            if slide_index < 0 or slide_index >= len(prs.slides):
                print(f"Warnung: Folie {slide_number} existiert nicht in {file}. Übersprungen.")
                continue

            # Hole die Folie und füge sie zur neuen Präsentation hinzu
            slide = prs.slides[slide_index]
            copy_slide(new_presentation, slide)

    # Speichere die neue Präsentation
    new_presentation.save(output_file)
    print(f"Neue Präsentation gespeichert unter: {output_file}")


def copy_slide(presentation, slide):
    """
    Kopiert eine Folie aus einer Präsentation in eine andere.

    :param presentation: Ziel-Präsentation.
    :param slide: Die Folie, die kopiert werden soll.
    """
    # Füge eine leere Folie mit einem Standardlayout hinzu
    layout = presentation.slide_layouts[0]  # Verwende ein Standardlayout
    new_slide = presentation.slides.add_slide(layout)

    # Kopiere Inhalte der Folie
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        if shape.has_text_frame:
            new_shape = new_slide.shapes.add_shape(
                shape.shape_type, shape.left, shape.top, shape.width, shape.height
            )
            new_shape.text = shape.text_frame.text
        elif shape.shape_type == 13:  # Bild
            cache_image(new_slide, shape)
        elif shape.shape_type == 14:  # Hintergrundbild (Background)
            cache_image(new_slide, shape)
        elif shape.shape_type == 1:  # AutoShape (z.B. Ovale)
            new_shape = new_slide.shapes.add_shape(
                shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height
            )
            if shape.has_text_frame:
                new_shape.text = shape.text_frame.text
        elif shape.shape_type == 19:  # Placeholder
            for placeholder in new_slide.placeholders:
                if placeholder.placeholder_format.idx == shape.placeholder_format.idx:
                    if shape.has_text_frame:
                        placeholder.text = shape.text_frame.text


def cache_image(new_slide, shape):
    with NamedTemporaryFile(delete=False, suffix=".png") as tmpfile:
        tmpfile.write(shape.image.blob)
        tmpfile.flush()
        new_slide.shapes.add_picture(
            tmpfile.name, shape.left, shape.top, shape.width, shape.height
        )
    os.remove(tmpfile.name)


if __name__ == "__main__":
    # Beispiel: Extrahieren von Folien aus verschiedenen Präsentationen
    slide_ranges = [
        ("../../Source/Präsentation_LF02-v2.pptx", [1, 2, 3]),
        ("../../Source/Topologie-Präsentation-v2.pptx", [2, 3])
    ]
    output_file = "../../Outcome/merged_presentation.pptx"
    extract_and_merge_slides(slide_ranges, output_file)

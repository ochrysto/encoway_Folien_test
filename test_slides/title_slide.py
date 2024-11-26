import datetime
from pptx import Presentation

# Data input
template = "../../Material/encoway PowerPoint-12-2023-Master.pptx"
title_content = input("Titel eingeben ")
subtitle_name = input("Name eingeben ")


# get date
date = datetime.datetime.now()
subtitle_date = f"{date.day}.{date.month}.{date.year}"

# open presentation
prs = Presentation(template)

# choose title layout and add title slide
title_slide_layout = prs.slide_layouts[0]
title_slide = prs.slides.add_slide(title_slide_layout)

# change text
title = title_slide.shapes.title
title.text = title_content

# change subtitle text
subtitle = title_slide.placeholders[10]
subtitle.text = f"{subtitle_name} | {subtitle_date}"

# save presentation
prs.save("test.pptx")

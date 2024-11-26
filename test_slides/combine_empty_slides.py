from pptx import Presentation
from pptx.util import Inches


class SlideNine:
    def __init__(self, title, subtitle, shorttext, longtext, img_url):
        self.title = title
        self.subtitle = subtitle
        self.shorttext = shorttext
        self.longtext = longtext
        self.img_url = img_url


data1 = SlideNine(
    "Machine",
    "feugiat a, tellus",
    "Lorem ipsum dolor sit amet, consectetuer adipiscing elit",
    "Aenean commodo ligula eget dolor. Aenean massa. Cum sociis natoque penatibus et magnis dis parturient montes, nascetur ridiculus mus. Donec quam felis, ultricies nec, pellentesque eu, pretium quis, sem. Nulla consequat massa quis enim. Donec pede justo, fringilla vel",
    "../../Material/images/pexels-pixabay-48799.jpg")
data2 = SlideNine(
    "Aenean imperdiet",
    "Etiam ultricies nisi vel augue",
    "sit amet adipiscing sem neque",
    "Aliquet nec, vulputate eget, arcu. In enim justo, rhoncus ut, imperdiet a, venenatis vitae, justo. Nullam dictum felis eu pede mollis pretium. Integer tincidunt. Cras dapibus. Vivamus elementum semper nisi. Aenean vulputate eleifend tellus. Aenean leo ligula, porttitor eu, consequat vitae, eleifend ac, enim. Aliquam lorem ante, dapibus in, viverra quis, feugiat a, tellus. Phasellus viverra nulla ut",
    "../../Material/images/pexels-pixabay-159201.jpg")
data3 = SlideNine(
    "Celeifend tellus",
    "Curabitur ullamcorper ultricies nisi. Nam eget dui. Etiam rhoncus",
    "Lorem ipsum dolor sit amet, consectetuer adipiscing elit",
    " Quisque rutrum. Aenean imperdiet. Etiam ultricies nisi vel augue. Curabitur ullamcorper ultricies nisi. Nam eget dui. Etiam rhoncus. Maecenas tempus, tellus eget condimentum rhoncus, sem quam semper libero, sit amet adipiscing sem neque",
    "../../Material/images/pexels-shvetsa-5953723.jpg")

data = [data1, data2, data3]
title_headline = input("Titel der Präsentation: ")
#number_of_slides = input("Anzahl der Folien: ")
file_name = input("Dateiname(ohne Dateiendung): ")


def create_title_slide(prs, title_text):
    slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(slide_layout)
    title = title_slide.shapes.title
    title.text = title_text


def create_slide_9(prs, title, subtitle, shorttext, longtext, img_url):
    slide_layout = prs.slide_layouts[9]
    slide_9 = prs.slides.add_slide(slide_layout)
    title_frame = slide_9.shapes.title
    title_frame.text = title
    subtitle_frame = slide_9.shapes.placeholders[0]
    subtitle_frame.text = subtitle
    text_box = slide_9.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    text_frame = text_box.text_frame
    text_frame.text = shorttext
    longtext_paragraph = text_frame.add_paragraph()
    longtext_paragraph.text = longtext
    slide_9.shapes.add_picture(img_url, left=Inches(5.5), top=Inches(2), height=Inches(5))


prs = Presentation()

create_title_slide(prs, title_headline)

for date in data:
    create_slide_9(prs, date.title, date.subtitle, date.shorttext, date.longtext, date.img_url)

prs.save(f"{file_name} .pptx")

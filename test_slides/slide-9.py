from pptx import Presentation
from pptx.util import Inches

template = "../../Material/encoway PowerPoint-12-2023-Master.pptx"
#img_path = "../../Material/images/pexels-pixabay-48799.jpg"
# input
headline = input("Überschrift: ")
subheader = input("Untertitel")
bullet_point_1 = input("Bulletpoint 1: ")
bullet_point_2 = input("Bulletpoint 2: ")
bullet_point_3 = input("Bulletpoint 3: ")
bullet_point_4 = input("Bulletpoint 4: ")
img_path = input("Pfad zum Bild: ")

prs = Presentation(template)

slide_9_layout = prs.slide_layouts[9]
slide_9 = prs.slides.add_slide(slide_9_layout)

shapes = slide_9.shapes

title = shapes.placeholders[0]
title.text = headline

subtitle = shapes.placeholders[36]
subtitle.text = subheader

#bullet points
body_left = shapes.placeholders[11]
bullet_points = body_left.text_frame
bullet_points.text = bullet_point_1

paragraph_two = bullet_points.add_paragraph()
paragraph_two.text = bullet_point_2

paragraph_three = bullet_points.add_paragraph()
paragraph_three.text = bullet_point_3

paragraph_four = bullet_points.add_paragraph()
paragraph_four.text = bullet_point_4

# subitem
# paragraph_one_subitem = text_frame.add_paragraph()
# paragraph_one_subitem.text = "Unterpunkt"
# paragraph_one_subitem.level = 1

# img
left = Inches(6.85)
top = Inches(1.75)
height = Inches(5)
img = shapes.add_picture(img_path, left, top, height=height)

prs.save("slide-9-test.pptx")